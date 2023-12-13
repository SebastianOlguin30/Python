from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Crear una presentación
prs = Presentation()

# Títulos de las diapositivas
titles = [
    "Mantén tus Dispositivos Actualizados",
    "Utiliza Contraseñas Fuertes",
    "Implementa la Autenticación de Dos Factores (2FA)",
    "Ten Cuidado con los Correos Electrónicos Phishing",
    "Actualiza tu Conocimiento sobre Phishing",
    "Respeta la Privacidad en Redes Sociales",
    "Protege tus Dispositivos Móviles",
    "Cuidado con las Redes Wi-Fi Públicas",
    "No Desactive las Actualizaciones Automáticas",
    "Haz Uso Responsable de los USB y Dispositivos Externos"
]

# Contenido de las diapositivas
contents = [
    "Asegúrate de que tu sistema operativo, software antivirus y otras aplicaciones estén siempre actualizados con las últimas correcciones de seguridad.",
    "Crea contraseñas únicas y fuertes para cada cuenta. Incluye letras, números y caracteres especiales, y evita contraseñas obvias como '123456' o 'password'.",
    "Habilita la autenticación de dos factores siempre que sea posible. Esto proporciona una capa adicional de seguridad, incluso si alguien conoce tu contraseña.",
    "Sé escéptico ante los correos electrónicos no solicitados. No hagas clic en enlaces ni descargues archivos adjuntos de remitentes desconocidos. Verifica la legitimidad de los correos electrónicos antes de proporcionar información sensible.",
    "Mantente informado sobre las tácticas de phishing más recientes. Los ataques de phishing evolucionan constantemente, y estar alerta es clave para evitar caer en trampas.",
    "Ajusta la configuración de privacidad en tus cuentas de redes sociales para limitar la cantidad de información personal que compartes públicamente.",
    "Utiliza contraseñas o códigos PIN en tus dispositivos móviles. Habilita la función de bloqueo remoto en caso de pérdida o robo.",
    "Evita realizar transacciones financieras o acceder a información confidencial en redes Wi-Fi públicas. Utiliza una red privada virtual (VPN) para mayor seguridad.",
    "Mantén activadas las actualizaciones automáticas en tu sistema operativo y software para recibir las últimas correcciones de seguridad.",
    "Evita conectar dispositivos USB desconocidos a tu computadora. Estos pueden contener malware."
]

for title, content in zip(titles, contents):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Usando un layout con título y contenido

    # Establecer el título
    title_box = slide.shapes.title
    title_box.text = title

    # Establecer el contenido
    content_box = slide.placeholders[1]
    content_box.text = content

    # Formatear el título
    for paragraph in title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)  # Color azul oscuro

    # Alinear y formatear el contenido
    for paragraph in content_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Guardar la presentación
prs.save('Presentacion_Seguridad_Informatica.pptx')